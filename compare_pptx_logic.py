"""
compare_pptx_logic.py

Core comparison logic extracted and simplified for integration with Flask.
Produces a color-coded HTML report (no CSV required for the web UI).
Uses deep_translator for back-translation to English for similarity checks.
"""

import json, time, html
from pathlib import Path
from pptx import Presentation
from deep_translator import GoogleTranslator
from difflib import SequenceMatcher
from datetime import datetime

def extract_text_blocks(pptx_path):
    prs = Presentation(pptx_path)
    blocks = []
    for s_idx, slide in enumerate(prs.slides, start=1):
        for sh_idx, shape in enumerate(slide.shapes):
            if not (hasattr(shape, "text_frame") and shape.has_text_frame):
                continue
            tf = shape.text_frame
            for p_idx, paragraph in enumerate(tf.paragraphs):
                raw = paragraph.text
                text = "" if raw is None else raw.strip()
                blocks.append({
                    "slide": s_idx,
                    "shape_index": sh_idx,
                    "para_index": p_idx,
                    "text": text
                })
    return blocks

def similarity_score(a, b):
    a = (a or "").strip().lower()
    b = (b or "").strip().lower()
    if not a and not b:
        return 1.0
    return SequenceMatcher(None, a, b).ratio()

def load_cache(path: Path):
    if path.exists():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}

def save_cache(cache: dict, path: Path):
    path.write_text(json.dumps(cache, ensure_ascii=False, indent=2), encoding="utf-8")

def build_html_report(rows, html_path: Path, eng_name: str, trans_name: str):
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    style = """
    <style>
      body { font-family: Arial, Helvetica, sans-serif; margin: 18px; color:#222; }
      h1 { text-align:center; color:#1f4e79; }
      table { border-collapse: collapse; width:100%; }
      th, td { border:1px solid #ddd; padding:8px; vertical-align:top; }
      th { background:#333; color:#fff; position: sticky; top:0; z-index:2; }
      td.small { width:90px; font-size:13px; text-align:center; }
      tr.ok { background:#e6f4ea; }
      tr.review { background:#fff7e6; }
      tr.missing { background:#fbeaea; }
      .mono { font-family: "Courier New", monospace; font-size:12px; color:#333; }
      .caption { margin-bottom:12px; color:#333; }
      .legend { margin:8px 0 18px 0; }
      .legend span { display:inline-block; padding:6px 10px; margin-right:10px; border-radius:4px; }
    </style>
    """
    legend = (
        '<div class="legend">'
        '<span style="background:#e6f4ea">OK (>= threshold)</span>'
        '<span style="background:#fff7e6">Needs Review (&lt; threshold)</span>'
        '<span style="background:#fbeaea">Missing / Extra</span>'
        '</div>'
    )
    header = f"<h1>Translation Comparison — {html.escape(eng_name)} ↔ {html.escape(trans_name)}</h1>"
    meta = f"<div class='caption'>Generated: {now}</div>"

    rows_html = []
    rows_html.append("<table>")
    rows_html.append("<thead><tr>"
                     "<th class='small'>Slide</th>"
                     "<th class='small'>Shape</th>"
                     "<th class='small'>Para</th>"
                     "<th>English text</th>"
                     "<th>Translated text</th>"
                     "<th>Back-translated (to EN)</th>"
                     "<th class='small'>Sim (%)</th>"
                     "<th class='small'>Status</th>"
                     "</tr></thead><tbody>")

    for r in rows:
        cls = "ok"
        status_lower = (r["status"] or "").lower()
        if "needs" in status_lower or "review" in status_lower:
            cls = "review"
        if "missing" in status_lower or "extra" in status_lower:
            cls = "missing"

        rows_html.append(
            "<tr class='%s'>%s</tr>" % (
                cls,
                "".join([
                    f"<td class='small'>{r['slide']}</td>",
                    f"<td class='small'>{r['shape_index']}</td>",
                    f"<td class='small'>{r['para_index']}</td>",
                    "<td>" + html.escape(r['english_text']) + "</td>",
                    "<td>" + html.escape(r['translated_text']) + "</td>",
                    "<td class='mono'>" + html.escape(r.get('back_translated','')) + "</td>",
                    f"<td class='small'>{r.get('similarity',0)}</td>",
                    f"<td class='small'>{html.escape(r['status'])}</td>",
                ])
            )
        )

    rows_html.append("</tbody></table>")
    html_content = f"<!doctype html><html><head><meta charset='utf-8'><title>Translation Comparison</title>{style}</head><body>{header}{meta}{legend}{''.join(rows_html)}</body></html>"

    html_path.write_text(html_content, encoding="utf-8")

def compare_presentations(eng_path: Path, trans_path: Path, outdir: Path, cache_file: Path, threshold: float=0.70, backtrans_target: str="en"):
    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)
    html_file = outdir / "translation_comparison_report.html"

    en_blocks = extract_text_blocks(eng_path)
    tr_blocks = extract_text_blocks(trans_path)

    en_map = { (b["slide"], b["shape_index"], b["para_index"]): b["text"] for b in en_blocks }
    tr_map = { (b["slide"], b["shape_index"], b["para_index"]): b["text"] for b in tr_blocks }

    all_keys = sorted(set(list(en_map.keys()) + list(tr_map.keys())), key=lambda k: (k[0], k[1], k[2]))

    cache = load_cache(cache_file)
    translator = GoogleTranslator(source="auto", target=backtrans_target)

    rows = []
    for key in all_keys:
        slide, shape_idx, para_idx = key
        en_text = en_map.get(key, "")
        tr_text = tr_map.get(key, "")

        if not tr_text:
            rows.append({
                "slide": slide,
                "shape_index": shape_idx,
                "para_index": para_idx,
                "english_text": en_text,
                "translated_text": tr_text,
                "back_translated": "",
                "similarity": 0.0,
                "status": "Missing translation" if en_text else "Empty block"
            })
            continue

        cache_key = f"bt::{tr_text}"
        if cache_key in cache:
            back_translated = cache[cache_key]
        else:
            try:
                back_translated = translator.translate(tr_text)
            except Exception:
                back_translated = ""
            cache[cache_key] = back_translated
            time.sleep(0.18)

        sim = similarity_score(en_text, back_translated)
        status = "OK" if sim >= threshold else "Needs Review"
        if not en_text:
            status = "Extra translation (no source)"
        rows.append({
            "slide": slide,
            "shape_index": shape_idx,
            "para_index": para_idx,
            "english_text": en_text,
            "translated_text": tr_text,
            "back_translated": back_translated,
            "similarity": round(sim * 100, 2),
            "status": status
        })

    save_cache(cache, cache_file)
    build_html_report(rows, html_file, eng_path.name, trans_path.name)
    return None, html_file

def compare_presentations_with_progress(
    eng_path, trans_path, outdir, cache_file, threshold=0.7, backtrans_target="en", progress=None
):
    from compare_pptx_logic import (
    extract_text_blocks,
    load_cache,
    save_cache,
    build_html_report,
    similarity_score,
)


    if progress is None:
        progress = {"current": 0, "total": 0, "done": False, "message": ""}

    en_blocks = extract_text_blocks(eng_path)
    tr_blocks = extract_text_blocks(trans_path)
    en_map = {(b["slide"], b["shape_index"], b["para_index"]): b["text"] for b in en_blocks}
    tr_map = {(b["slide"], b["shape_index"], b["para_index"]): b["text"] for b in tr_blocks}
    all_keys = sorted(set(en_map.keys()) | set(tr_map.keys()), key=lambda k: (k[0], k[1], k[2]))

    progress["total"] = len(all_keys)
    cache = load_cache(cache_file)
    translator = GoogleTranslator(source="auto", target=backtrans_target)
    rows = []

    for i, key in enumerate(all_keys, start=1):
        slide, shape_idx, para_idx = key
        progress["current"] = i
        progress["message"] = f"Processing slide {slide} (item {i}/{len(all_keys)})"

        en_text = en_map.get(key, "")
        tr_text = tr_map.get(key, "")
        if not tr_text:
            rows.append(
                {
                    "slide": slide,
                    "shape_index": shape_idx,
                    "para_index": para_idx,
                    "english_text": en_text,
                    "translated_text": "",
                    "back_translated": "",
                    "similarity": 0,
                    "status": "Missing translation" if en_text else "Empty block",
                }
            )
            continue

        cache_key = f"bt::{tr_text}"
        back_translated = cache.get(cache_key, "")
        if not back_translated:
            try:
                back_translated = translator.translate(tr_text)
                cache[cache_key] = back_translated
            except Exception:
                back_translated = ""
            save_cache(cache, cache_file)

        sim = similarity_score(en_text, back_translated)
        status = "OK" if sim >= threshold else "Needs Review"
        if not en_text:
            status = "Extra translation (no source)"
        rows.append(
            {
                "slide": slide,
                "shape_index": shape_idx,
                "para_index": para_idx,
                "english_text": en_text,
                "translated_text": tr_text,
                "back_translated": back_translated,
                "similarity": round(sim * 100, 2),
                "status": status,
            }
        )

    html_file = outdir / "translation_comparison_report.html"
    build_html_report(rows, html_file, eng_path.name, trans_path.name)
    progress["done"] = True
    progress["message"] = "✅ Report ready"

    return None, html_file

